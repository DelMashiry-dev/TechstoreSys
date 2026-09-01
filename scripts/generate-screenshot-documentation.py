#!/usr/bin/env python3
"""Capture TECHSTORESys module screenshots and build illustrated PDF documentation."""

from __future__ import annotations

import argparse
import importlib.util
import json
import subprocess
import sys
import time
import urllib.error
import urllib.request
from datetime import date
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
SHOT_DIR = ROOT / "docs" / "screenshots"
MANIFEST_PATH = SHOT_DIR / "manifest.json"
OUT_ROOT = ROOT / "TECHSTORESys-System-Documentation-Screenshots.pdf"
OUT_APP = ROOT / "app" / "TECHSTORESys-System-Documentation-Screenshots.pdf"
OUT_PORTABLE = ROOT / "dist" / "TECHSTORES-Portable" / "app" / "TECHSTORESys-System-Documentation-Screenshots.pdf"
PPT_ROOT = ROOT / "TECHSTORESys-System-Presentation.pptx"
PPT_APP = ROOT / "app" / "TECHSTORESys-System-Presentation.pptx"
PPT_PORTABLE = ROOT / "dist" / "TECHSTORES-Portable" / "app" / "TECHSTORESys-System-Presentation.pptx"

BASE_URL = "http://127.0.0.1:8080/app/"
LOGIN_USER = "admin"
LOGIN_PASS = "admin123"
DOC_VERSION = "2.2"

# Role accounts: preset key, username, password, display label
ROLE_ACCOUNTS = [
    ("gate", "rp", "rp123", "RP Gate"),
    ("storeman", "storeman", "storeman123", "Storeman"),
    ("workshop", "workshop", "workshop123", "Workshop"),
    ("orderly", "orderly", "orderly123", "Orderly Room"),
    ("comms", "sysadmin", "sysadmin123", "Department comms (SysAdmin)"),
    ("dp", "dp", "dp123", "Directorate Procurement (DP)"),
    ("gs", "gsdesk", "gsdesk123", "GS Branch"),
    ("daf", "daf", "daf123", "DAF / MANAC"),
    ("aiad", "aiad", "aiad123", "AIAD Due Diligence"),
    ("supplier", "nixzimo", "nixzimo123", "Supplier (Nixzimo)"),
    ("store_officer", "store", "store123", "Store Officer (full ledgers)"),
    ("viewer", "viewer", "view123", "Viewer (read-only)"),
]

# Import labels / module list from the text documentation generator.
_spec = importlib.util.spec_from_file_location(
    "doc_pdf",
    ROOT / "scripts" / "generate-system-documentation-pdf.py",
)
_doc_mod = importlib.util.module_from_spec(_spec)
assert _spec and _spec.loader
_spec.loader.exec_module(_doc_mod)
MODULE_LABELS: dict[str, str] = _doc_mod.MODULE_LABELS
MODULE_IDS: list[str] = _doc_mod.MODULE_IDS

_spec_desc = importlib.util.spec_from_file_location(
    "screen_desc",
    ROOT / "scripts" / "screen-descriptions.py",
)
_desc_mod = importlib.util.module_from_spec(_spec_desc)
assert _spec_desc and _spec_desc.loader
_spec_desc.loader.exec_module(_desc_mod)
get_screen_description = _desc_mod.get_screen_description

STAKEHOLDER_CAPTURES = [
    ("stakeholder-desk-dp", "stakeholder-desk", {"stkDesk": "dp"}, "DP Window - procurement portal"),
    ("stakeholder-desk-gs", "stakeholder-desk", {"stkDesk": "gs"}, "GS Branch Window - F1 endorsement"),
    (
        "stakeholder-desk-daf-procurement",
        "stakeholder-desk",
        {"stkDesk": "daf"},
        "DAF Window - MANAC / supplier payment",
    ),
    (
        "stakeholder-desk-daf-creditors",
        "stakeholder-desk",
        {"stkDesk": "daf", "stkDafTab": "creditors"},
        "DAF Window - Creditors import & chase",
    ),
    ("stakeholder-desk-aiad", "stakeholder-desk", {"stkDesk": "aiad"}, "Due Diligence Window - AIAD"),
    ("stakeholder-desk-supplier", "stakeholder-desk", {"stkDesk": "supplier"}, "Supplier Window"),
]

def role_home_slugs() -> list[str]:
    return [f"role-home-{key}" for key, _, _, _ in ROLE_ACCOUNTS]


def role_nav_slugs() -> list[str]:
    return [f"role-nav-{key}" for key, _, _, _ in ROLE_ACCOUNTS]


SECTIONS = [
    ("Login & shell", ["00-login", "01-app-shell"]),
    ("Role-based home dashboards", role_home_slugs()),
    ("Role-based navigation (sidebar)", role_nav_slugs()),
    (
        "Dashboard & communications",
        ["dashboard", "orderly-room", "it-dir-comms", "portals-board"],
    ),
    (
        "IT Dir department desks",
        [
            "dept-sysadmin",
            "dept-workshop",
            "dept-compengr",
            "dept-swengr",
            "dept-ictsec",
            "dept-itts",
            "dept-admin",
            "dept-gate",
        ],
    ),
    (
        "Stakeholder portals",
        [
            "stakeholder-desk-dp",
            "stakeholder-desk-gs",
            "stakeholder-desk-daf-procurement",
            "stakeholder-desk-daf-creditors",
            "stakeholder-desk-aiad",
            "stakeholder-desk-supplier",
        ],
    ),
    (
        "General ledger & stores",
        [
            "gl-2200600002",
            "gl-2200600003",
            "gl-220200002",
            "gl-2201900002",
            "gl-3112210001",
            "voucher-module",
            "stock-take",
            "unit-checks",
            "accommodation-stores",
            "financial-year-bids",
        ],
    ),
    (
        "Equipment, loans & returns",
        [
            "unit-equipment",
            "ict-accountability",
            "ict-distribution",
            "temporary-loans",
            "permanent-loans",
            "monthly-returns",
            "techstores-equipment-register",
        ],
    ),
    (
        "Procurement & creditors",
        [
            "unit-requisitions",
            "doc-import",
            "dp-f1-form",
            "cost-comparative-schedule",
            "spec-evaluation",
            "guide-quotation",
            "dp-procurement",
            "purchase-orders",
            "delivery-note",
            "undelivered-orders",
            "supplier-debts",
            "suppliers-contracts",
        ],
    ),
    (
        "Workshop & compare tools",
        [
            "workshop-repairs",
            "workshop-receipt-cert",
            "laptop-compare",
            "ict-compare",
            "gate-register",
        ],
    ),
    (
        "ZNA Q forms index",
        ["zna-q-forms-index"],
    ),
    (
        "ZNA Q & SVCS forms",
        [
            "zna-q-982",
            "zna-q-178",
            "zna-q-1033",
            "zna-q-1043",
            "zna-q-80",
            "zna-svcs-890",
            "zna-q-1179",
            "zna-q-987",
            "zna-q-3977",
            "zna-q-1157",
            "zna-q-985",
            "zna-q-1",
            "zna-q-998",
            "zna-q-1680",
            "zna-q-3",
            "zna-q-31",
            "zna-q-40",
            "zna-q-1049",
            "zna-q-1229",
            "zna-q-1571",
            "zna-q-1954",
            "zna-svcs-1045",
        ],
    ),
    (
        "Administration & help",
        [
            "duties-roles",
            "process-guides",
            "system-help",
            "reports-module",
            "user-management",
            "release-cut",
        ],
    ),
]


def ascii_safe(text: str) -> str:
    if not text:
        return ""
    return (
        text.replace("\u2014", "-")
        .replace("\u2013", "-")
        .replace("\u2019", "'")
        .replace("\u2018", "'")
        .replace("\u2022", "-")
        .encode("latin-1", "replace")
        .decode("latin-1")
    )


def build_capture_targets() -> list[dict]:
    targets: list[dict] = [
        {
            "slug": "00-login",
            "label": "Login screen",
            "kind": "login",
            "module_id": None,
            "options": {},
        },
        {
            "slug": "01-app-shell",
            "label": "Application shell (sidebar + header)",
            "kind": "shell",
            "module_id": "dashboard",
            "options": {},
        },
    ]
    skip = {"stakeholder-desk"}
    for mid in MODULE_IDS:
        if mid in skip:
            continue
        targets.append(
            {
                "slug": mid,
                "label": MODULE_LABELS.get(mid, mid),
                "kind": "module",
                "module_id": mid,
                "options": {},
            }
        )
    for slug, module_id, options, label in STAKEHOLDER_CAPTURES:
        targets.append(
            {
                "slug": slug,
                "label": label,
                "kind": "module",
                "module_id": module_id,
                "options": options,
            }
        )
    return targets


def ordered_slugs() -> list[str]:
    ordered: list[str] = []
    seen: set[str] = set()
    for _, slugs in SECTIONS:
        for slug in slugs:
            if slug not in seen:
                ordered.append(slug)
                seen.add(slug)
    for target in build_capture_targets():
        slug = target["slug"]
        if slug not in seen:
            ordered.append(slug)
            seen.add(slug)
    return ordered


def wait_for_server(url: str = "http://127.0.0.1:8080/api/health", timeout: float = 45.0) -> None:
    deadline = time.time() + timeout
    while time.time() < deadline:
        try:
            with urllib.request.urlopen(url, timeout=2) as resp:
                if resp.status == 200:
                    return
        except (urllib.error.URLError, TimeoutError, OSError):
            time.sleep(0.5)
    raise RuntimeError(f"Server not reachable at {url}. Start with: python server.py")


def login_user(page, username: str, password: str) -> None:
    if page.locator("#loginUsername").count() == 0 or not page.locator("#loginUsername").is_visible():
        page.goto(BASE_URL, wait_until="networkidle", timeout=120000)
    page.fill("#loginUsername", username)
    page.fill("#loginPassword", password)
    page.click("#loginForm button[type='submit']")
    page.wait_for_function(
        "() => !document.body.classList.contains('app-locked')",
        timeout=120000,
    )
    page.wait_for_timeout(1200)


def logout_user(page) -> None:
    btn = page.locator("#logoutBtn")
    if btn.count() and btn.is_visible():
        btn.click()
        page.wait_for_function(
            "() => document.body.classList.contains('app-locked')",
            timeout=30000,
        )
        page.wait_for_timeout(400)


def capture_role_screenshots(page, manifest: dict, force: bool = False) -> None:
    for key, username, password, label in ROLE_ACCOUNTS:
        home_slug = f"role-home-{key}"
        nav_slug = f"role-nav-{key}"
        home_path = SHOT_DIR / f"{home_slug}.jpg"
        nav_path = SHOT_DIR / f"{nav_slug}.jpg"
        need_home = force or not home_path.exists()
        need_nav = force or not nav_path.exists()
        if not need_home and not need_nav:
            manifest["screenshots"][home_slug] = {
                "file": home_path.name,
                "label": f"{label} - home dashboard",
                "login": username,
                "ok": True,
                "skipped": True,
            }
            manifest["screenshots"][nav_slug] = {
                "file": nav_path.name,
                "label": f"{label} - sidebar navigation",
                "login": username,
                "ok": True,
                "skipped": True,
            }
            continue

        logout_user(page)
        login_user(page, username, password)
        page.evaluate(
            """async () => {
                if (typeof navigateToModule === 'function') {
                    await navigateToModule('dashboard');
                }
            }"""
        )
        page.wait_for_timeout(900)

        if need_home:
            try:
                page.locator(".main-content").screenshot(path=str(home_path), type="jpeg", quality=88)
                ok = home_path.exists()
                err = ""
            except Exception as exc:  # noqa: BLE001
                ok = False
                err = str(exc)
                print(f"[warn] role home failed for {home_slug}: {exc}", file=sys.stderr)
            manifest["screenshots"][home_slug] = {
                "file": home_path.name if home_path.exists() else "",
                "label": f"{label} - home dashboard",
                "login": username,
                "ok": ok,
                "error": err if not ok else "",
            }
            print(f"{'OK' if ok else 'FAIL'}  {home_slug}")

        if need_nav:
            try:
                page.locator(".container.app-shell").screenshot(path=str(nav_path), type="jpeg", quality=88)
                ok = nav_path.exists()
                err = ""
            except Exception as exc:  # noqa: BLE001
                ok = False
                err = str(exc)
                print(f"[warn] role nav failed for {nav_slug}: {exc}", file=sys.stderr)
            manifest["screenshots"][nav_slug] = {
                "file": nav_path.name if nav_path.exists() else "",
                "label": f"{label} - sidebar navigation",
                "login": username,
                "ok": ok,
                "error": err if not ok else "",
            }
            print(f"{'OK' if ok else 'FAIL'}  {nav_slug}")


def maybe_start_server() -> subprocess.Popen | None:
    try:
        wait_for_server(timeout=2)
        return None
    except RuntimeError:
        pass
    proc = subprocess.Popen(
        [sys.executable, str(ROOT / "server.py")],
        cwd=str(ROOT),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    wait_for_server(timeout=60)
    return proc


def capture_screenshots(force: bool = False, roles_only: bool = False) -> dict:
    from playwright.sync_api import sync_playwright

    SHOT_DIR.mkdir(parents=True, exist_ok=True)
    targets = build_capture_targets()
    manifest: dict = load_manifest() if roles_only else {"screenshots": {}}
    manifest["generated"] = date.today().isoformat()
    manifest["base_url"] = BASE_URL

    server_proc = maybe_start_server()
    try:
        with sync_playwright() as pw:
            browser = pw.chromium.launch(headless=True)
            context = browser.new_context(viewport={"width": 1440, "height": 900}, device_scale_factor=1)
            page = context.new_page()

            if not roles_only:
                # Login screen
                login_path = SHOT_DIR / "00-login.jpg"
                if force or not login_path.exists():
                    page.goto(BASE_URL, wait_until="networkidle", timeout=120000)
                    page.wait_for_selector("#loginScreen", timeout=30000)
                    page.locator("#loginScreen").screenshot(path=str(login_path), type="jpeg", quality=88)
                manifest["screenshots"]["00-login"] = {
                    "file": login_path.name,
                    "label": "Login screen",
                    "ok": login_path.exists(),
                }

                login_user(page, LOGIN_USER, LOGIN_PASS)

            for target in ([] if roles_only else targets):
                slug = target["slug"]
                out_path = SHOT_DIR / f"{slug}.jpg"
                if out_path.exists() and not force and slug != "01-app-shell":
                    manifest["screenshots"][slug] = {
                        "file": out_path.name,
                        "label": target["label"],
                        "ok": True,
                        "skipped": True,
                    }
                    continue

                ok = True
                err = ""
                try:
                    if target["kind"] == "login":
                        continue
                    if target["kind"] == "shell":
                        page.locator(".container.app-shell").screenshot(
                            path=str(out_path), type="jpeg", quality=88
                        )
                    else:
                        module_id = target["module_id"]
                        options = target.get("options") or {}
                        page.evaluate(
                            """async ({ moduleId, options }) => {
                                if (typeof navigateToModule !== 'function') {
                                    throw new Error('navigateToModule missing');
                                }
                                await navigateToModule(moduleId, options);
                            }""",
                            {"moduleId": module_id, "options": options},
                        )
                        page.wait_for_timeout(1200)
                        page.wait_for_function(
                            """(moduleId) => {
                                const el = document.getElementById(moduleId);
                                return el && el.style.display !== 'none';
                            }""",
                            arg=module_id,
                            timeout=30000,
                        )
                        page.wait_for_timeout(500)
                        page.locator(".main-content").screenshot(
                            path=str(out_path), type="jpeg", quality=88
                        )
                except Exception as exc:  # noqa: BLE001
                    ok = False
                    err = str(exc)
                    print(f"[warn] capture failed for {slug}: {exc}", file=sys.stderr)

                manifest["screenshots"][slug] = {
                    "file": out_path.name if out_path.exists() else "",
                    "label": target["label"],
                    "module_id": target.get("module_id"),
                    "ok": ok and out_path.exists(),
                    "error": err,
                }
                print(f"{'OK' if manifest['screenshots'][slug]['ok'] else 'FAIL'}  {slug}")

            if roles_only:
                page.goto(BASE_URL, wait_until="networkidle", timeout=120000)
            capture_role_screenshots(page, manifest, force=force)

            browser.close()
    finally:
        if server_proc and server_proc.poll() is None:
            server_proc.terminate()
            try:
                server_proc.wait(timeout=5)
            except subprocess.TimeoutExpired:
                server_proc.kill()

    MANIFEST_PATH.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    return manifest


def load_manifest() -> dict:
    if MANIFEST_PATH.exists():
        return json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    return {"screenshots": {}}


def fit_image_to_page(pdf, img_path: Path, max_w: float = 190.0, max_h: float = 235.0) -> None:
    try:
        from PIL import Image
    except ImportError:
        pdf.image(str(img_path), w=max_w)
        return

    with Image.open(img_path) as im:
        px_w, px_h = im.size
    if px_w <= 0 or px_h <= 0:
        return
    aspect = px_h / px_w
    w = max_w
    h = w * aspect
    if h > max_h:
        h = max_h
        w = h / aspect
    pdf.image(str(img_path), w=w, h=h)


class ScreenshotPDF:
    def __init__(self) -> None:
        from fpdf import FPDF

        class _PDF(FPDF):
            def header(self):
                if self.page_no() == 1:
                    return
                self.set_font("Helvetica", "B", 9)
                self.set_text_color(27, 94, 59)
                self.cell(0, 6, "TECHSTORESys - Illustrated System Documentation", new_x="LMARGIN", new_y="NEXT")
                self.set_font("Helvetica", "", 8)
                self.set_text_color(110, 110, 110)
                self.cell(
                    0,
                    4,
                    f"Zimbabwe National Army | v{DOC_VERSION} | {date.today().strftime('%d %B %Y')}",
                    new_x="LMARGIN",
                    new_y="NEXT",
                )
                self.ln(2)

            def footer(self):
                self.set_y(-12)
                self.set_font("Helvetica", "I", 8)
                self.set_text_color(120, 120, 120)
                self.cell(0, 8, f"Page {self.page_no()}/{{nb}}", align="C")

        self.pdf = _PDF(orientation="P", unit="mm", format="A4")
        self.pdf.alias_nb_pages()
        self.pdf.set_auto_page_break(auto=True, margin=12)

    def cover(self, shot_count: int) -> None:
        pdf = self.pdf
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 22)
        pdf.set_text_color(27, 94, 59)
        pdf.cell(0, 14, "TECHSTORESys", new_x="LMARGIN", new_y="NEXT", align="C")
        pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 10, "Illustrated System Documentation", new_x="LMARGIN", new_y="NEXT", align="C")
        pdf.ln(4)
        pdf.set_font("Helvetica", "", 11)
        pdf.set_text_color(60, 60, 60)
        pdf.multi_cell(
            0,
            6,
            ascii_safe(
                "IT Directorate Tech Stores - General Ledger, stores ledgers, procurement portals, "
                "ZNA Q forms, workshop tools, and role-based workspaces."
            ),
            align="C",
        )
        pdf.ln(6)
        pdf.set_font("Helvetica", "", 10)
        pdf.cell(0, 6, f"Document version {DOC_VERSION}", new_x="LMARGIN", new_y="NEXT", align="C")
        pdf.cell(0, 6, date.today().strftime("%d %B %Y"), new_x="LMARGIN", new_y="NEXT", align="C")
        pdf.cell(0, 6, f"{shot_count} illustrated screens with descriptions", new_x="LMARGIN", new_y="NEXT", align="C")
        pdf.cell(0, 6, "Companion PowerPoint: TECHSTORESys-System-Presentation.pptx", new_x="LMARGIN", new_y="NEXT", align="C")
        pdf.ln(8)
        pdf.set_font("Helvetica", "I", 9)
        pdf.multi_cell(
            0,
            5,
            "Screens captured from the live application at 1440x900. "
            "Includes admin module walkthrough plus role-based home dashboards and sidebars. "
            "Regenerate: python scripts/generate-screenshot-documentation.py",
            align="C",
        )

    def toc(self, entries: list[tuple[str, str]]) -> None:
        pdf = self.pdf
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 14)
        pdf.set_text_color(27, 94, 59)
        pdf.cell(0, 10, "Contents", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(2)
        pdf.set_font("Helvetica", "", 9)
        pdf.set_text_color(40, 40, 40)
        for section, _slug in entries:
            pdf.cell(0, 5, f"  {ascii_safe(section)}", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(4)

    def section_title(self, title: str) -> None:
        pdf = self.pdf
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 16)
        pdf.set_text_color(27, 94, 59)
        pdf.cell(0, 12, ascii_safe(title), new_x="LMARGIN", new_y="NEXT")
        pdf.set_draw_color(184, 207, 196)
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(4)

    def screen_page(
        self,
        title: str,
        module_id: str | None,
        img_path: Path | None,
        login_user: str | None = None,
        description: dict[str, str] | None = None,
    ) -> None:
        pdf = self.pdf
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 12)
        pdf.set_text_color(30, 30, 30)
        pdf.multi_cell(0, 6, ascii_safe(title))
        if module_id:
            pdf.set_font("Helvetica", "", 8)
            pdf.set_text_color(100, 100, 100)
            pdf.cell(0, 5, ascii_safe(f"Module ID: {module_id}"), new_x="LMARGIN", new_y="NEXT")
        elif login_user:
            pdf.set_font("Helvetica", "", 8)
            pdf.set_text_color(100, 100, 100)
            pdf.cell(0, 5, ascii_safe(f"Login: {login_user}"), new_x="LMARGIN", new_y="NEXT")
        if description:
            pdf.ln(1)
            pdf.set_font("Helvetica", "B", 9)
            pdf.set_text_color(27, 94, 59)
            pdf.cell(0, 5, "What it is", new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("Helvetica", "", 8.5)
            pdf.set_text_color(50, 50, 50)
            pdf.multi_cell(0, 4.2, ascii_safe(description.get("what", "")))
            pdf.ln(0.5)
            pdf.set_font("Helvetica", "B", 9)
            pdf.set_text_color(27, 94, 59)
            pdf.cell(0, 5, "How it works", new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("Helvetica", "", 8.5)
            pdf.multi_cell(0, 4.2, ascii_safe(description.get("how", "")))
        pdf.ln(2)
        img_max_h = 175.0 if description else 235.0
        if img_path and img_path.exists():
            fit_image_to_page(pdf, img_path, max_h=img_max_h)
        else:
            pdf.set_font("Helvetica", "I", 10)
            pdf.set_text_color(160, 40, 40)
            pdf.cell(0, 8, "[Screenshot not available]", new_x="LMARGIN", new_y="NEXT")

    def write(self, path: Path) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        self.pdf.output(str(path))


def slug_label(slug: str, manifest: dict) -> str:
    info = manifest.get("screenshots", {}).get(slug, {})
    if info.get("label"):
        return info["label"]
    return MODULE_LABELS.get(slug, slug.replace("-", " ").title())


def generate_pdf(manifest: dict | None = None) -> None:
    manifest = manifest or load_manifest()
    shots = manifest.get("screenshots", {})
    doc = ScreenshotPDF()

    slugs = ordered_slugs()
    present = [s for s in slugs if shots.get(s, {}).get("ok")]
    doc.cover(len(present))

    toc_entries: list[tuple[str, str]] = []
    for section_title, section_slugs in SECTIONS:
        for slug in section_slugs:
            if shots.get(slug, {}).get("ok"):
                toc_entries.append((f"{ascii_safe(section_title)} - {ascii_safe(slug_label(slug, manifest))}", slug))
    doc.toc(toc_entries)

    current_section = None
    for section_title, section_slugs in SECTIONS:
        section_has = any(shots.get(s, {}).get("ok") for s in section_slugs)
        if not section_has:
            continue
        doc.section_title(section_title)
        for slug in section_slugs:
            info = shots.get(slug, {})
            if not info.get("ok"):
                continue
            img_name = info.get("file") or f"{slug}.jpg"
            img_path = SHOT_DIR / img_name
            module_id = info.get("module_id")
            if not module_id and slug.startswith("role-"):
                module_id = None
            elif not module_id and slug not in ("00-login", "01-app-shell") and not slug.startswith("role-"):
                module_id = slug
            doc.screen_page(
                slug_label(slug, manifest),
                module_id,
                img_path,
                login_user=info.get("login"),
                description=get_screen_description(slug),
            )

    for out in (OUT_ROOT, OUT_APP, OUT_PORTABLE):
        doc.write(out)
        print(f"Wrote {out}")


def _ppt_add_title_box(slide, text: str, left, top, width, height, size: int = 28, bold: bool = True, color=None):
    from pptx.util import Pt

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    return box


def _ppt_add_bullets(slide, lines: list[tuple[str, str]], left, top, width, height):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for label, body in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"{label}: {body}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x34, 0x40, 0x54)
        p.space_after = Pt(6)
    return box


def generate_presentation(manifest: dict | None = None) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    manifest = manifest or load_manifest()
    shots = manifest.get("screenshots", {})
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    green = RGBColor(0x1B, 0x5E, 0x3B)
    white = RGBColor(0xFF, 0xFF, 0xFF)

    slugs = ordered_slugs()
    present_count = sum(1 for s in slugs if shots.get(s, {}).get("ok"))

    # Title slide
    slide = prs.slides.add_slide(blank)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = green
    bar.line.fill.background()
    _ppt_add_title_box(
        slide, "TECHSTORESys", Inches(0.6), Inches(0.25), Inches(12), Inches(0.7), size=36, color=white
    )
    _ppt_add_title_box(
        slide,
        "IT Directorate Tech Stores - System Overview",
        Inches(0.6),
        Inches(1.8),
        Inches(12),
        Inches(0.6),
        size=22,
        bold=False,
    )
    _ppt_add_title_box(
        slide,
        f"Zimbabwe National Army | v{DOC_VERSION} | {date.today().strftime('%d %B %Y')}",
        Inches(0.6),
        Inches(2.5),
        Inches(12),
        Inches(0.5),
        size=14,
        bold=False,
    )
    _ppt_add_title_box(
        slide,
        f"{present_count} screens - stores, GL, procurement portals, ZNA Q forms, role workspaces",
        Inches(0.6),
        Inches(3.2),
        Inches(12),
        Inches(0.8),
        size=13,
        bold=False,
    )

    # Agenda slide
    slide = prs.slides.add_slide(blank)
    _ppt_add_title_box(slide, "Contents", Inches(0.5), Inches(0.35), Inches(12), Inches(0.6), size=24, color=green)
    agenda_lines = [
        title for title, section_slugs in SECTIONS if any(shots.get(s, {}).get("ok") for s in section_slugs)
    ]
    y = 1.1
    for i, title in enumerate(agenda_lines, 1):
        _ppt_add_title_box(
            slide,
            f"{i}. {title}",
            Inches(0.7),
            Inches(y),
            Inches(11.5),
            Inches(0.35),
            size=14,
            bold=False,
        )
        y += 0.38
        if y > 6.8:
            break

    for section_title, section_slugs in SECTIONS:
        if not any(shots.get(s, {}).get("ok") for s in section_slugs):
            continue

        # Section divider
        slide = prs.slides.add_slide(blank)
        bar = slide.shapes.add_shape(1, Inches(0), Inches(2.8), prs.slide_width, Inches(1.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = green
        bar.line.fill.background()
        _ppt_add_title_box(
            slide, section_title, Inches(0.6), Inches(3.0), Inches(12), Inches(1.2), size=30, color=white
        )

        for slug in section_slugs:
            info = shots.get(slug, {})
            if not info.get("ok"):
                continue
            img_name = info.get("file") or f"{slug}.jpg"
            img_path = SHOT_DIR / img_name
            title = slug_label(slug, manifest)
            desc = get_screen_description(slug)
            module_id = info.get("module_id")
            meta = f"Module: {module_id}" if module_id else (f"Login: {info.get('login')}" if info.get("login") else "")

            slide = prs.slides.add_slide(blank)
            _ppt_add_title_box(slide, title, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.55), size=18, color=green)
            if meta:
                _ppt_add_title_box(
                    slide, meta, Inches(0.4), Inches(0.72), Inches(12), Inches(0.3), size=10, bold=False
                )
            _ppt_add_bullets(
                slide,
                [("What it is", desc.get("what", "")), ("How it works", desc.get("how", ""))],
                Inches(0.4),
                Inches(1.05),
                Inches(12.4),
                Inches(1.35),
            )
            if img_path.exists():
                pic_w = Inches(12.4)
                pic_left = Inches(0.45)
                pic_top = Inches(2.55)
                pic = slide.shapes.add_picture(str(img_path), pic_left, pic_top, width=pic_w)
                max_h = Inches(4.75)
                if pic.height > max_h:
                    ratio = max_h / pic.height
                    pic.width = int(pic.width * ratio)
                    pic.height = int(pic.height * ratio)
                    pic.left = pic_left + (pic_w - pic.width) // 2

    for out in (PPT_ROOT, PPT_APP, PPT_PORTABLE):
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
        print(f"Wrote {out}")


def main() -> int:
    parser = argparse.ArgumentParser(description="Capture module screenshots and build illustrated PDF.")
    parser.add_argument("--capture-only", action="store_true", help="Only capture screenshots.")
    parser.add_argument("--pdf-only", action="store_true", help="Only build PDF from existing screenshots.")
    parser.add_argument("--ppt-only", action="store_true", help="Only build PowerPoint from existing screenshots.")
    parser.add_argument("--force", action="store_true", help="Re-capture even if JPEG exists.")
    parser.add_argument(
        "--roles-only",
        action="store_true",
        help="Only capture role home/nav screens (merge into existing manifest).",
    )
    args = parser.parse_args()

    if args.pdf_only or args.ppt_only:
        manifest = load_manifest()
        if not manifest.get("screenshots"):
            print("No manifest found. Run capture first.", file=sys.stderr)
            return 1
        if args.ppt_only:
            generate_presentation(manifest)
        else:
            generate_pdf(manifest)
            generate_presentation(manifest)
        return 0

    manifest = capture_screenshots(force=args.force, roles_only=args.roles_only)
    if not args.capture_only:
        generate_pdf(manifest)
        generate_presentation(manifest)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
